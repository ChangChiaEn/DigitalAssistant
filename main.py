"""
Digital Assistant - Desktop AI Assistant
Python + PyWebView + Colab LLM Backend
"""
import webview
import json
import os
import subprocess
import platform
import threading
import datetime
import tempfile

# ===== System Skills API (exposed to JavaScript) =====

class AssistantAPI:
    """Python backend exposed to the web UI via pywebview."""

    WORKSPACE = os.path.join(os.path.expanduser('~'), 'Desktop', 'AssistantOutput')

    def __init__(self):
        self._settings_path = os.path.join(os.path.dirname(__file__), 'settings.json')
        self._settings = self._load_settings()
        os.makedirs(self.WORKSPACE, exist_ok=True)

    # --- Settings ---
    def _load_settings(self):
        try:
            with open(self._settings_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except:
            return {}

    def _save_settings(self):
        with open(self._settings_path, 'w', encoding='utf-8') as f:
            json.dump(self._settings, f, indent=2, ensure_ascii=False)

    def get_setting(self, key):
        return self._settings.get(key, '')

    def set_setting(self, key, value):
        self._settings[key] = value
        self._save_settings()

    # --- Launch App ---
    def launch_app(self, app_name):
        app_map = {
            'notepad': 'notepad.exe',
            '記事本': 'notepad.exe',
            'calculator': 'calc.exe',
            '計算機': 'calc.exe',
            'browser': 'start https://www.google.com',
            '瀏覽器': 'start https://www.google.com',
            'explorer': 'explorer.exe',
            '檔案總管': 'explorer.exe',
            'cmd': 'cmd.exe',
            'terminal': 'wt.exe',
            '終端機': 'wt.exe',
            'vscode': 'code',
            'spotify': 'start spotify:',
            'discord': 'start discord:',
        }
        cmd = app_map.get(app_name.lower(), app_name)
        try:
            subprocess.Popen(f'start "" "{cmd}"', shell=True)
            return json.dumps({'success': True, 'message': f'已啟動: {app_name}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Open URL ---
    def open_url(self, url):
        try:
            os.startfile(url)
            return json.dumps({'success': True, 'message': f'已開啟: {url}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Open File/Folder ---
    def open_path(self, path):
        try:
            os.startfile(path)
            return json.dumps({'success': True, 'message': f'已開啟: {path}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Run Command ---
    def run_command(self, command):
        try:
            result = subprocess.run(
                command, shell=True, capture_output=True,
                text=True, timeout=15, encoding='utf-8', errors='replace'
            )
            output = result.stdout.strip() or result.stderr.strip()
            return json.dumps({'success': result.returncode == 0, 'message': output})
        except subprocess.TimeoutExpired:
            return json.dumps({'success': False, 'message': '指令逾時 (15秒)'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- System Info ---
    def system_info(self):
        mem = os.popen('wmic OS get FreePhysicalMemory,TotalVisibleMemorySize /Value').read()
        info = {
            'platform': platform.system(),
            'hostname': platform.node(),
            'cpus': os.cpu_count(),
            'version': platform.version(),
        }
        for line in mem.strip().split('\n'):
            line = line.strip()
            if line.startswith('TotalVisibleMemorySize='):
                total_kb = int(line.split('=')[1])
                info['totalMemory'] = f'{total_kb / 1048576:.1f} GB'
            elif line.startswith('FreePhysicalMemory='):
                free_kb = int(line.split('=')[1])
                info['freeMemory'] = f'{free_kb / 1048576:.1f} GB'
        return json.dumps({'success': True, 'message': json.dumps(info, ensure_ascii=False)})

    # --- Clipboard ---
    def clipboard_read(self):
        try:
            result = subprocess.run(
                'powershell -command "Get-Clipboard"',
                shell=True, capture_output=True, text=True, encoding='utf-8', errors='replace'
            )
            return json.dumps({'success': True, 'message': result.stdout.strip()})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    def clipboard_write(self, text):
        try:
            process = subprocess.Popen(
                'powershell -command "Set-Clipboard -Value $input"',
                shell=True, stdin=subprocess.PIPE, encoding='utf-8'
            )
            process.communicate(input=text)
            return json.dumps({'success': True, 'message': '已複製到剪貼簿'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # ===== NEW SKILLS =====

    # --- Create PPT ---
    def create_ppt(self, title, slides_json, theme='dark'):
        """Create a styled PowerPoint file. slides_json = [{"title":"...", "content":"..."}]
        theme: dark, corporate, nature, warm, ocean, minimal"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.dml.color import RGBColor

            slides = json.loads(slides_json) if isinstance(slides_json, str) else slides_json

            # Theme palettes
            THEMES = {
                'dark': {
                    'bg': (0x1A, 0x1A, 0x2E), 'card': (0x16, 0x21, 0x3E),
                    'accent': (0x53, 0x3C, 0xD4), 'accent2': (0x7C, 0x6A, 0xF0),
                    'text': (0xFF, 0xFF, 0xFF), 'sub': (0xA0, 0xA0, 0xB8), 'light': (0xE0, 0xE0, 0xEE),
                },
                'corporate': {
                    'bg': (0xFF, 0xFF, 0xFF), 'card': (0xF0, 0xF2, 0xF5),
                    'accent': (0x00, 0x52, 0xCC), 'accent2': (0x33, 0x7A, 0xFF),
                    'text': (0x1A, 0x1A, 0x2E), 'sub': (0x66, 0x66, 0x80), 'light': (0x33, 0x33, 0x44),
                },
                'nature': {
                    'bg': (0xF5, 0xF0, 0xEB), 'card': (0xE8, 0xE0, 0xD5),
                    'accent': (0x2D, 0x6A, 0x4F), 'accent2': (0x52, 0xB7, 0x88),
                    'text': (0x2D, 0x2D, 0x2D), 'sub': (0x6B, 0x70, 0x5C), 'light': (0x3D, 0x3D, 0x3D),
                },
                'warm': {
                    'bg': (0x2D, 0x1B, 0x1B), 'card': (0x3D, 0x25, 0x25),
                    'accent': (0xE8, 0x6C, 0x3A), 'accent2': (0xFF, 0x9A, 0x6C),
                    'text': (0xFF, 0xFF, 0xFF), 'sub': (0xB0, 0x98, 0x90), 'light': (0xE8, 0xDB, 0xD5),
                },
                'ocean': {
                    'bg': (0x0B, 0x1D, 0x33), 'card': (0x12, 0x2A, 0x45),
                    'accent': (0x00, 0xB4, 0xD8), 'accent2': (0x48, 0xCA, 0xE4),
                    'text': (0xFF, 0xFF, 0xFF), 'sub': (0x90, 0xBC, 0xD0), 'light': (0xCA, 0xE9, 0xF5),
                },
                'minimal': {
                    'bg': (0xFF, 0xFF, 0xFF), 'card': (0xF8, 0xF8, 0xF8),
                    'accent': (0x22, 0x22, 0x22), 'accent2': (0x55, 0x55, 0x55),
                    'text': (0x11, 0x11, 0x11), 'sub': (0x88, 0x88, 0x88), 'light': (0x33, 0x33, 0x33),
                },
            }
            t = THEMES.get(theme, THEMES['dark'])
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Color palette from theme
            BG_DARK = RGBColor(*t['bg'])
            BG_CARD = RGBColor(*t['card'])
            ACCENT = RGBColor(*t['accent'])
            ACCENT_LIGHT = RGBColor(*t['accent2'])
            WHITE = RGBColor(*t['text'])
            GRAY = RGBColor(*t['sub'])
            LIGHT = RGBColor(*t['light'])

            def add_bg(slide, color):
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = color

            def add_shape(slide, left, top, w, h, fill_color, radius=None):
                shape = slide.shapes.add_shape(
                    1, left, top, w, h  # MSO_SHAPE.RECTANGLE
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_color
                shape.line.fill.background()
                return shape

            def set_text(tf, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
                tf.clear()
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(size)
                p.font.color.rgb = color
                p.font.bold = bold
                p.font.name = 'Microsoft JhengHei'
                p.alignment = align

            # ── Title Slide ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            add_bg(slide, BG_DARK)

            # Accent bar left
            add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

            # Title
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Microsoft JhengHei'

            # Subtitle / date
            txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.0), Inches(10), Inches(0.6))
            set_text(txBox2.text_frame,
                     datetime.datetime.now().strftime('%Y . %m . %d'),
                     size=16, color=GRAY)

            # Decorative line
            add_shape(slide, Inches(1.5), Inches(3.8), Inches(3), Pt(3), ACCENT_LIGHT)

            # ── Content Slides ──
            for idx, s in enumerate(slides):
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
                add_bg(slide, BG_DARK)

                # Top accent bar
                add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT)

                # Slide number
                txNum = slide.shapes.add_textbox(Inches(11.8), Inches(0.3), Inches(1), Inches(0.5))
                set_text(txNum.text_frame, f'{idx + 1:02d}', size=14, color=GRAY, align=PP_ALIGN.RIGHT)

                # Section title
                txTitle = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1.0))
                set_text(txTitle.text_frame, s.get('title', ''), size=32, color=WHITE, bold=True)

                # Divider
                add_shape(slide, Inches(0.8), Inches(1.5), Inches(2), Pt(3), ACCENT_LIGHT)

                # Content card
                card = add_shape(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(5.0), BG_CARD)
                card.shadow.inherit = False

                # Content text with bullet points
                content_text = s.get('content', '')
                txContent = slide.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(11), Inches(4.5))
                tf = txContent.text_frame
                tf.word_wrap = True

                lines = [l.strip() for l in content_text.split('\n') if l.strip()]
                for i, line in enumerate(lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    # Add bullet marker if not already present
                    if not line.startswith(('-', '*', '>')):
                        p.text = f'  {line}'
                    else:
                        p.text = f'  {line.lstrip("-*> ")}'
                    p.font.size = Pt(18)
                    p.font.color.rgb = LIGHT
                    p.font.name = 'Microsoft JhengHei'
                    p.space_after = Pt(10)

            # ── End Slide ──
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_bg(slide, BG_DARK)
            add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(4), ACCENT)

            txEnd = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(9), Inches(1.5))
            set_text(txEnd.text_frame, 'Thank You', size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

            txSub = slide.shapes.add_textbox(Inches(2), Inches(4.3), Inches(9), Inches(0.6))
            set_text(txSub.text_frame, title, size=18, color=GRAY, align=PP_ALIGN.CENTER)

            filename = f"{title.replace(' ', '_')}_{datetime.datetime.now().strftime('%H%M%S')}.pptx"
            filepath = os.path.join(self.WORKSPACE, filename)
            prs.save(filepath)
            os.startfile(filepath)
            return json.dumps({'success': True, 'message': f'已建立並開啟: {filepath}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Create Word Document ---
    def create_docx(self, title, content):
        """Create a styled Word document."""
        try:
            from docx import Document
            from docx.shared import Inches, Pt, Cm, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.section import WD_ORIENT

            doc = Document()

            # Page setup
            section = doc.sections[0]
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(2.8)
            section.right_margin = Cm(2.8)

            # Style: Normal
            style_normal = doc.styles['Normal']
            style_normal.font.name = 'Microsoft JhengHei'
            style_normal.font.size = Pt(11)
            style_normal.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
            style_normal.paragraph_format.space_after = Pt(6)
            style_normal.paragraph_format.line_spacing = 1.5

            # Style: Title
            style_title = doc.styles['Title']
            style_title.font.name = 'Microsoft JhengHei'
            style_title.font.size = Pt(26)
            style_title.font.bold = True
            style_title.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            style_title.paragraph_format.space_after = Pt(4)

            # Style: Headings
            for level in range(1, 4):
                h_style = doc.styles[f'Heading {level}']
                h_style.font.name = 'Microsoft JhengHei'
                h_style.font.color.rgb = RGBColor(0x2D, 0x2D, 0x5E)
                h_style.font.bold = True
                sizes = {1: 20, 2: 16, 3: 13}
                h_style.font.size = Pt(sizes[level])
                h_style.paragraph_format.space_before = Pt(18 if level == 1 else 12)
                h_style.paragraph_format.space_after = Pt(8)

            # ── Cover section ──
            # Add some top spacing
            for _ in range(3):
                doc.add_paragraph('')

            # Title
            p_title = doc.add_paragraph(title, style='Title')
            p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Accent line
            p_line = doc.add_paragraph()
            p_line.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p_line.add_run('━' * 30)
            run.font.color.rgb = RGBColor(0x53, 0x3C, 0xD4)
            run.font.size = Pt(10)

            # Date
            p_date = doc.add_paragraph()
            p_date.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p_date.add_run(datetime.datetime.now().strftime('%Y-%m-%d'))
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(0x88, 0x88, 0x99)

            # Page break after cover
            doc.add_page_break()

            # ── Content ──
            lines = content.split('\n')
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue

                # Detect headings by markers
                if stripped.startswith('# '):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith('## '):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith('### '):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith(('-', '*', '>')):
                    # Bullet point
                    text = stripped.lstrip('-*> ').strip()
                    p = doc.add_paragraph(text, style='List Bullet')
                    p.paragraph_format.left_indent = Cm(1)
                else:
                    p = doc.add_paragraph(stripped)

            # ── Footer ──
            footer_section = doc.sections[0]
            footer = footer_section.footer
            footer_para = footer.paragraphs[0]
            footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = footer_para.add_run(f'{title}  |  Generated by Digital Assistant')
            run.font.size = Pt(8)
            run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)
            run.font.name = 'Microsoft JhengHei'

            filename = f"{title.replace(' ', '_')}_{datetime.datetime.now().strftime('%H%M%S')}.docx"
            filepath = os.path.join(self.WORKSPACE, filename)
            doc.save(filepath)
            os.startfile(filepath)
            return json.dumps({'success': True, 'message': f'已建立並開啟: {filepath}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Create Excel ---
    def create_xlsx(self, title, data_json):
        """Create an Excel file. data_json = [["col1","col2"],["val1","val2"]]"""
        try:
            from openpyxl import Workbook

            data = json.loads(data_json) if isinstance(data_json, str) else data_json
            wb = Workbook()
            ws = wb.active
            ws.title = title
            for row in data:
                ws.append(row)

            filename = f"{title.replace(' ', '_')}_{datetime.datetime.now().strftime('%H%M%S')}.xlsx"
            filepath = os.path.join(self.WORKSPACE, filename)
            wb.save(filepath)
            os.startfile(filepath)
            return json.dumps({'success': True, 'message': f'已建立並開啟: {filepath}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Write Text File ---
    def write_file(self, filename, content):
        """Write content to a text file in workspace."""
        try:
            filepath = os.path.join(self.WORKSPACE, filename)
            with open(filepath, 'w', encoding='utf-8') as f:
                f.write(content)
            return json.dumps({'success': True, 'message': f'已儲存: {filepath}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Read File ---
    def read_file(self, filepath):
        """Read content from a file."""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(10000)  # max 10k chars
            return json.dumps({'success': True, 'message': content})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- List Files ---
    def list_files(self, directory):
        """List files in a directory."""
        try:
            path = directory or self.WORKSPACE
            items = []
            for name in os.listdir(path)[:50]:
                full = os.path.join(path, name)
                is_dir = os.path.isdir(full)
                size = os.path.getsize(full) if not is_dir else 0
                items.append(f"{'[DIR]' if is_dir else f'{size:>8}'} {name}")
            return json.dumps({'success': True, 'message': '\n'.join(items) or '(empty)'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Screenshot ---
    def take_screenshot(self):
        """Take a screenshot and save to workspace."""
        try:
            from PIL import ImageGrab
            img = ImageGrab.grab()
            filename = f"screenshot_{datetime.datetime.now().strftime('%H%M%S')}.png"
            filepath = os.path.join(self.WORKSPACE, filename)
            img.save(filepath)
            return json.dumps({'success': True, 'message': f'已截圖: {filepath}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Get Date/Time ---
    def get_datetime(self):
        """Get current date and time."""
        now = datetime.datetime.now()
        return json.dumps({'success': True, 'message': json.dumps({
            'date': now.strftime('%Y-%m-%d'),
            'time': now.strftime('%H:%M:%S'),
            'weekday': ['一', '二', '三', '四', '五', '六', '日'][now.weekday()],
            'full': now.strftime('%Y年%m月%d日 %H:%M 星期') + ['一', '二', '三', '四', '五', '六', '日'][now.weekday()]
        }, ensure_ascii=False)})

    # --- Search Web (open browser with search) ---
    def search_web(self, query):
        """Search the web by opening browser."""
        try:
            import urllib.parse
            url = f"https://www.google.com/search?q={urllib.parse.quote(query)}"
            os.startfile(url)
            return json.dumps({'success': True, 'message': f'已搜尋: {query}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Fetch News / Web Search Results ---
    def fetch_news(self, query, max_results='5'):
        """Search the web and return actual content summaries using DuckDuckGo."""
        try:
            from duckduckgo_search import DDGS
            max_r = int(max_results)
            results = []
            
            # Use query as-is, let AI decide what to search
            # Only do minimal cleanup: remove extra spaces
            clean_query = ' '.join(query.split()).strip()

            with DDGS() as ddgs:
                seen_titles = set()
                for search_q in search_queries:
                    for r in ddgs.text(search_q, max_results=max_r):
                        title = r.get('title', '')
                        # Skip duplicates
                        if title in seen_titles:
                            continue
                        seen_titles.add(title)
                        results.append({
                            'title': r.get('title', ''),
                            'body': r.get('body', ''),
                            'href': r.get('href', ''),
                        })
                        if len(results) >= max_r:
                            break
                    if len(results) >= max_r:
                        break

            if not results:
                return json.dumps({'success': False, 'message': f'找不到關於「{query}」的結果'})

            # Format as readable text for AI to process
            formatted = []
            for i, r in enumerate(results, 1):
                formatted.append(f"[{i}] {r['title']}\n{r['body']}\n來源: {r['href']}")

            summary = '\n\n'.join(formatted)
            return json.dumps({
                'success': True,
                'message': summary,
                'results': results,
                'query': query
            }, ensure_ascii=False)
        except Exception as e:
            return json.dumps({'success': False, 'message': f'搜尋失敗: {str(e)}'})

    # --- Kill Process ---
    def kill_process(self, process_name):
        """Kill a running process by name."""
        try:
            result = subprocess.run(
                f'taskkill /f /im {process_name}',
                shell=True, capture_output=True, text=True, encoding='utf-8', errors='replace'
            )
            msg = result.stdout.strip() or result.stderr.strip()
            return json.dumps({'success': result.returncode == 0, 'message': msg})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Set Volume (Windows) ---
    def set_volume(self, level):
        """Set system volume (0-100)."""
        try:
            # Use PowerShell to set volume
            ps_cmd = f'''
            $wshShell = New-Object -ComObject WScript.Shell;
            1..50 | ForEach-Object {{ $wshShell.SendKeys([char]174) }};
            1..{int(int(level)/2)} | ForEach-Object {{ $wshShell.SendKeys([char]175) }}
            '''
            subprocess.run(f'powershell -command "{ps_cmd}"', shell=True, timeout=10)
            return json.dumps({'success': True, 'message': f'音量已設定為 {level}%'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # --- Notify (Windows Toast) ---
    def notify(self, title, message):
        """Show a Windows notification."""
        try:
            ps_cmd = f'''
            [Windows.UI.Notifications.ToastNotificationManager, Windows.UI.Notifications, ContentType = WindowsRuntime] > $null
            $template = [Windows.UI.Notifications.ToastNotificationManager]::GetTemplateContent([Windows.UI.Notifications.ToastTemplateType]::ToastText02)
            $textNodes = $template.GetElementsByTagName("text")
            $textNodes.Item(0).AppendChild($template.CreateTextNode("{title}")) > $null
            $textNodes.Item(1).AppendChild($template.CreateTextNode("{message}")) > $null
            $toast = [Windows.UI.Notifications.ToastNotification]::new($template)
            [Windows.UI.Notifications.ToastNotificationManager]::CreateToastNotifier("Digital Assistant").Show($toast)
            '''
            subprocess.run(f'powershell -command "{ps_cmd}"', shell=True, timeout=10)
            return json.dumps({'success': True, 'message': f'已發送通知: {title}'})
        except Exception as e:
            return json.dumps({'success': False, 'message': str(e)})

    # ===== Speech Recognition (Python-side) =====
    def start_listening(self, lang):
        """Record from microphone and return recognized text using Google Speech API."""
        try:
            import speech_recognition as sr
            recognizer = sr.Recognizer()
            recognizer.energy_threshold = 300
            recognizer.dynamic_energy_threshold = True
            recognizer.pause_threshold = 1.5      # wait 1.5s of silence before stopping
            recognizer.phrase_threshold = 0.3
            recognizer.non_speaking_duration = 0.8

            with sr.Microphone() as source:
                recognizer.adjust_for_ambient_noise(source, duration=0.5)
                audio = recognizer.listen(source, timeout=10, phrase_time_limit=30)

            text = recognizer.recognize_google(audio, language=lang)
            return json.dumps({'success': True, 'text': text})
        except Exception as e:
            error_msg = str(e)
            if 'timed out' in error_msg.lower():
                return json.dumps({'success': False, 'text': '', 'error': '未偵測到語音'})
            return json.dumps({'success': False, 'text': '', 'error': error_msg})

    # ===== AI Chat (calls Colab) =====
    def chat_with_ai(self, messages_json, skills_json):
        import requests

        api_url = self._settings.get('apiUrl', '')
        if not api_url:
            return json.dumps({'text': '[Error] 尚未設定 Colab API URL，請點擊齒輪設定。'})

        try:
            messages = json.loads(messages_json)
            skills = json.loads(skills_json)

            response = requests.post(
                f'{api_url}/chat',
                json={'messages': messages, 'skills': skills},
                timeout=30,
                headers={'Content-Type': 'application/json'}
            )

            if response.ok:
                return json.dumps(response.json(), ensure_ascii=False)
            else:
                return json.dumps({'text': f'[Error] API: {response.status_code}'})

        except requests.exceptions.Timeout:
            return json.dumps({'text': '[Error] AI 回應逾時，請確認 Colab 是否仍在運行。'})
        except requests.exceptions.ConnectionError:
            return json.dumps({'text': '[Error] 無法連線到 AI 後端，請確認 Colab 是否已啟動。'})
        except Exception as e:
            return json.dumps({'text': f'[Error] {str(e)}'})

    # --- Health Check ---
    def check_health(self):
        import requests
        api_url = self._settings.get('apiUrl', '')
        if not api_url:
            return json.dumps({'connected': False, 'message': '未設定'})

        try:
            response = requests.get(f'{api_url}/health', timeout=5)
            if response.ok:
                data = response.json()
                return json.dumps({
                    'connected': True,
                    'message': f"已連線 - {data.get('model', 'AI Ready')}"
                })
        except:
            pass

        return json.dumps({'connected': False, 'message': '未連線 - 請啟動 Colab'})

    # --- Window Controls ---
    def close_window(self):
        """Close the application window."""
        import webview
        for w in webview.windows:
            w.destroy()

    def minimize_window(self):
        """Minimize the application window."""
        import webview
        for w in webview.windows:
            w.minimize()


# ===== Launch =====
if __name__ == '__main__':
    api = AssistantAPI()

    window = webview.create_window(
        'Assistant',
        url=os.path.join(os.path.dirname(__file__), 'public', 'index.html'),
        js_api=api,
        width=420,
        height=680,
        resizable=True,
        on_top=True,
        frameless=True,
        easy_drag=True,
        background_color='#0e0f14',
    )

    webview.start(debug='--dev' in os.sys.argv)
